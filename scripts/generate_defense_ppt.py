# -*- coding: utf-8 -*-
"""Generate a concise graduation-defense PPT for the Smart Travel System."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "defense-presentation-smart-travel-final.pptx"
OUTLINE_PATH = ROOT / "defense-presentation-smart-travel-outline.md"
COVER_IMAGE = ROOT / "travel-system-uniapp" / "src" / "static" / "images" / "bg.jpg"

STUDENT_NAME = "【请填写姓名】"
STUDENT_CLASS = "【请填写班级】"
SCHOOL_NAME = "【学校名称】"
TITLE = "基于 Spring Boot 与 UniApp 的智能旅游系统设计与实现"

WIDE = Inches(13.333)
HIGH = Inches(7.5)

BLUE = RGBColor(25, 90, 170)
BLUE_DARK = RGBColor(18, 48, 96)
BLUE_LIGHT = RGBColor(225, 238, 255)
CYAN = RGBColor(14, 165, 233)
GREEN = RGBColor(16, 145, 110)
GREEN_LIGHT = RGBColor(222, 247, 236)
YELLOW = RGBColor(245, 178, 35)
ORANGE = RGBColor(240, 118, 69)
RED = RGBColor(220, 80, 80)
GRAY_900 = RGBColor(30, 41, 59)
GRAY_700 = RGBColor(51, 65, 85)
GRAY_500 = RGBColor(100, 116, 139)
GRAY_200 = RGBColor(226, 232, 240)
GRAY_100 = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)


SLIDES = [
    {
        "title": "毕业设计题目",
        "layout": "封面采用旅游背景图叠加蓝色半透明遮罩，题目居中突出，基本信息置于下方。",
        "visual": "封面背景图、题名、学生姓名、班级、学校名称、8 分钟答辩提示。",
        "speech": "各位老师好，我的毕业设计题目是《基于 Spring Boot 与 UniApp 的智能旅游系统设计与实现》。本系统面向游客出行前的信息获取、行程规划和出行后的内容分享场景，后续我将从技术栈、核心算法、系统架构、功能模块、数据库设计和成果展示六个方面进行汇报。",
    },
    {
        "title": "项目开发环境与技术栈",
        "layout": "五类技术卡片横向排布，下方用运行链路说明三端协作方式。",
        "visual": "技术栈图标式卡片：UniApp/Vue3、Spring Boot/MyBatis、MySQL/Redis、IDEA/HBuilderX/微信开发者工具、JDK/Node/小程序运行环境。",
        "speech": "系统采用前后端分离架构。用户端是 UniApp 和 Vue3 开发的微信小程序，管理端使用 Vue3、Vite 和 Element Plus，后端基于 Spring Boot 2.7、MyBatis 和 RESTful API。数据层使用 MySQL 保存核心业务数据，Redis 用于推荐热点结果缓存，运行环境包括 JDK8、Node.js 和微信小程序运行环境。",
    },
    {
        "title": "核心算法（一）：智能行程规划",
        "layout": "左侧为输入信息，右侧为五步流程图，底部展示数据落库结构。",
        "visual": "流程图：用户输入 -> 标签匹配 -> 候选 POI 排序 -> 分天编排 -> AI 文案/规则兜底 -> 落库。",
        "speech": "第一个核心设计是智能行程规划。系统接收城市、天数、兴趣标签、预算、同行人以及用户必选景点和美食。首先根据 content_tag 筛选同城市候选 POI，若没有匹配结果则使用全量兜底；然后将用户必选项前置，再按热度和评分排序；最后按天生成行程明细，并可调用 DeepSeek 生成概要和早中晚说明，AI 不可用时自动使用规则文案兜底。",
    },
    {
        "title": "核心算法（二）：个性化推荐",
        "layout": "顶部展示推荐数据来源，中部为评分公式，底部为推荐输出。",
        "visual": "算法图：用户标签权重 + 行为数据 + 内容热度 -> 综合评分 -> Top N 推荐；配 Redis 30 分钟缓存。",
        "speech": "第二个核心设计是个性化推荐。系统从 user_tag 获取用户兴趣权重，同时结合浏览、点赞、收藏、评论等行为，以及线路、游记、景点、美食自身的热度指标。不同内容有不同权重，例如线路会考虑收藏数、浏览数和使用次数，游记会考虑点赞、收藏、浏览和评论。最终按综合得分排序返回 Top N，并对部分热门结果使用 Redis 缓存 30 分钟，降低数据库压力。",
    },
    {
        "title": "系统整体架构设计",
        "layout": "从用户角色到客户端、API、业务服务、数据与外部服务的分层架构图。",
        "visual": "系统架构图：游客/管理员 -> 小程序/管理后台 -> REST API -> user/content/route/recommend/travel/common -> MySQL/Redis/文件/AI。",
        "speech": "系统整体上分为用户端、管理端和后端服务。游客通过微信小程序进行浏览、规划、打卡和内容互动；管理员通过 Web 后台进行用户、内容、线路和运营配置管理。两端统一通过 REST API 访问 Spring Boot 后端，后端按 user、content、route、recommend、travel 和 common 模块组织业务逻辑，再分别访问 MySQL、Redis、上传文件目录和 AI 服务。",
    },
    {
        "title": "系统功能模块划分",
        "layout": "按用户端、管理端、后端能力三组模块展示，每个模块配一句功能说明。",
        "visual": "模块结构图：用户模块、内容模块、行程模块、推荐模块、社区互动模块、后台管理模块。",
        "speech": "功能模块围绕旅游全过程展开。用户端提供首页推荐、城市景点美食浏览、智能规划、游记社区、打卡足迹和个人中心；管理端提供数据总览、用户管理、城市景点美食管理、游记审核、评论打卡管理、路线管理、推荐策略和等级配置；后端负责接口编排、业务校验、推荐计算、AI 文案生成、文件上传和统一异常处理。",
    },
    {
        "title": "数据库设计",
        "layout": "用 ER 关系图展示核心表，右侧列出关键字段与设计说明。",
        "visual": "ER 图：user、tag、city、scenic_spot、food、travel_route、travel_route_day、travel_route_poi、travel_note、comment、checkin_record、user_behavior。",
        "speech": "数据库使用 MySQL 的 smart_travel 库，核心表按用户域、内容域、行程域、社区域和行为域划分。用户和标签通过 user_tag 建立多对多偏好关系；城市关联景点和美食，并通过 content_tag 与标签建立关系；行程采用 travel_route、travel_route_day、travel_route_poi 三级结构，方便展示多日行程；用户行为表记录浏览、点赞等行为，为推荐算法提供数据支撑。",
    },
    {
        "title": "项目成果展示",
        "layout": "左侧为三张系统界面截图式面板，右侧为成果亮点与现场演示流程，底部致谢。",
        "visual": "运行界面展示：小程序首页/智能规划/管理后台；亮点总结：AI 行程、推荐、UGC、后台运营。",
        "speech": "最后展示系统成果。系统已完成小程序端、管理后台和后端接口联调，支持首页推荐、智能行程生成、景点美食详情、游记发布互动、景点打卡，以及后台内容管理和审核。现场演示时可以按登录、选择兴趣、生成行程、查看详情、发布游记、后台审核的顺序进行。我的汇报到此结束，感谢各位老师聆听，敬请批评指正。",
    },
]


def set_font(run, size=16, bold=False, color=GRAY_700):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def fill(shape, color, line=GRAY_200, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.color.rgb = line


def add_text(slide, text, x, y, w, h, size=16, color=GRAY_700, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    set_font(p.runs[0], size=size, bold=bold, color=color)
    return box


def add_bullets(slide, items, x, y, w, h, size=14, color=GRAY_700, bullet=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.06)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT
        if bullet:
            p.text = "• " + item
        set_font(p.runs[0], size=size, color=color)
    return box


def add_header(slide, title, page):
    add_text(slide, title, 0.55, 0.25, 9.8, 0.45, size=23, color=BLUE_DARK, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.82), Inches(1.2), Inches(0.04))
    fill(line, BLUE, line=BLUE)
    add_text(slide, f"{page:02d}/08", 12.1, 6.95, 0.7, 0.25, size=9, color=GRAY_500, align=PP_ALIGN.RIGHT)


def add_footer_note(slide, text):
    add_text(slide, text, 0.6, 7.0, 8.5, 0.25, size=9, color=GRAY_500)


def add_card(slide, x, y, w, h, title, body=None, color=WHITE, accent=BLUE, icon=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shp, color, line=GRAY_200)
    if icon:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.18), Inches(0.48), Inches(0.48))
        fill(circ, accent, line=accent)
        add_text(slide, icon, x + 0.18, y + 0.27, 0.48, 0.16, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tx = x + 0.78
        tw = w - 0.95
    else:
        tx = x + 0.22
        tw = w - 0.42
    add_text(slide, title, tx, y + 0.2, tw, 0.32, size=14, color=GRAY_900, bold=True)
    if body:
        add_text(slide, body, x + 0.22, y + 0.66, w - 0.44, h - 0.76, size=11, color=GRAY_700)
    return shp


def connector(slide, x1, y1, x2, y2, color=BLUE, width=2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def pill(slide, text, x, y, w, color=BLUE_LIGHT, line=BLUE, txt=BLUE_DARK):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.42))
    fill(shp, color, line=line)
    add_text(slide, text, x + 0.05, y + 0.11, w - 0.1, 0.16, size=10.5, color=txt, bold=True, align=PP_ALIGN.CENTER)
    return shp


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    if COVER_IMAGE.exists():
        slide.shapes.add_picture(str(COVER_IMAGE), Inches(0), Inches(0), width=WIDE, height=HIGH)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), WIDE, HIGH)
    fill(overlay, BLUE_DARK, line=BLUE_DARK, transparency=20)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), WIDE, Inches(0.18))
    fill(band, CYAN, line=CYAN)

    add_text(slide, "毕业设计答辩", 0.9, 0.95, 3.2, 0.35, size=16, color=BLUE_LIGHT, bold=True)
    add_text(slide, TITLE, 0.9, 1.55, 8.9, 1.1, size=32, color=WHITE, bold=True)
    add_text(slide, "智能旅游系统 · 微信小程序 · 管理后台 · AI 行程规划", 0.92, 2.78, 8.2, 0.36, size=15, color=BLUE_LIGHT)

    info = [
        f"学生姓名：{STUDENT_NAME}",
        f"班级：{STUDENT_CLASS}",
        f"学校名称：{SCHOOL_NAME}",
        "自述时间：8 分钟（含系统演示）",
    ]
    y = 5.1
    for item in info:
        add_text(slide, item, 0.95, y, 6.6, 0.28, size=14, color=WHITE)
        y += 0.38
    add_text(slide, "01/08", 12.1, 6.95, 0.7, 0.25, size=9, color=BLUE_LIGHT, align=PP_ALIGN.RIGHT)


def slide_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GRAY_100
    add_header(slide, "项目开发环境与技术栈", 2)

    cards = [
        ("前端用户端", "UniApp\nVue 3\n微信小程序", "UNI", BLUE),
        ("管理后台", "Vue 3\nVite\nElement Plus", "VUE", GREEN),
        ("后端服务", "Spring Boot 2.7\nMyBatis\nRESTful API", "SB", BLUE_DARK),
        ("数据与缓存", "MySQL 8.0\nRedis\n上传文件目录", "DB", ORANGE),
        ("开发与运行", "IDEA / HBuilderX\n微信开发者工具\nJDK 8 / Node.js", "DEV", CYAN),
    ]
    x = 0.62
    for title, body, icon, color in cards:
        add_card(slide, x, 1.32, 2.32, 2.65, title, body, color=WHITE, accent=color, icon=icon)
        x += 2.48

    add_text(slide, "运行链路", 0.75, 4.45, 1.2, 0.25, size=14, color=BLUE_DARK, bold=True)
    steps = [
        ("微信小程序", 0.8, "用户浏览、规划、互动"),
        ("管理后台", 3.05, "内容维护、审核、运营"),
        ("Spring Boot API", 5.35, "业务接口与服务编排"),
        ("MySQL / Redis", 7.95, "持久化与热点缓存"),
        ("AI / 文件服务", 10.45, "行程文案与图片访问"),
    ]
    for label, x, sub in steps:
        add_card(slide, x, 5.05, 1.75, 0.86, label, sub, color=WHITE, accent=BLUE)
    for x in [2.55, 4.85, 7.45, 9.95]:
        connector(slide, x, 5.48, x + 0.42, 5.48, color=GRAY_500, width=1.4)
    add_footer_note(slide, "讲解重点：强调三端协同、前后端分离、MySQL+Redis 的数据支撑。")


def slide_algorithm_route(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, "核心算法（一）：智能行程规划", 3)

    add_card(
        slide, 0.72, 1.22, 3.0, 4.8, "输入信息",
        "城市 / 天数\n兴趣标签\n预算与同行人\n必选景点、美食\n每日自定义选择",
        color=BLUE_LIGHT, accent=BLUE, icon="IN",
    )
    flow = [
        ("标签筛选", "content_tag 匹配景点/美食"),
        ("候选兜底", "无匹配时回退城市全量 POI"),
        ("优先排序", "必选项前置 + 热度/评分"),
        ("分天编排", "早餐/景点/午晚餐按天入库"),
        ("AI 文案", "DeepSeek 生成概要；失败规则兜底"),
    ]
    x = 4.15
    y = 1.12
    for idx, (t, b) in enumerate(flow, start=1):
        add_card(slide, x, y, 2.2, 0.88, f"{idx}. {t}", b, color=WHITE, accent=BLUE)
        if idx < len(flow):
            connector(slide, x + 1.1, y + 0.88, x + 1.1, y + 1.18, color=BLUE, width=1.4)
        y += 1.18

    add_text(slide, "落库结构", 7.65, 1.22, 1.2, 0.25, size=14, color=BLUE_DARK, bold=True)
    add_card(slide, 7.65, 1.72, 1.72, 0.78, "travel_route", "线路主表", color=WHITE, accent=GREEN)
    add_card(slide, 9.9, 1.72, 1.78, 0.78, "route_day", "第 N 天", color=WHITE, accent=GREEN)
    add_card(slide, 9.9, 3.08, 1.78, 0.78, "route_poi", "景点/美食明细", color=WHITE, accent=GREEN)
    connector(slide, 9.37, 2.1, 9.86, 2.1, color=GREEN, width=1.5)
    connector(slide, 10.79, 2.5, 10.79, 3.03, color=GREEN, width=1.5)

    add_card(
        slide, 7.65, 4.45, 4.02, 1.12, "为什么这样设计",
        "用标签缩小候选范围，减少无关 POI；用必选项保证用户意图；用规则兜底保证 AI 不可用时流程仍可完成。",
        color=GREEN_LIGHT, accent=GREEN, icon="WHY",
    )
    add_footer_note(slide, "讲解重点：输入、筛选、排序、分天、AI/兜底五步即可，控制在 1 分钟内。")


def slide_algorithm_recommend(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GRAY_100
    add_header(slide, "核心算法（二）：个性化推荐", 4)

    sources = [
        ("用户兴趣", "user_tag.weight"),
        ("用户行为", "浏览 / 点赞 / 收藏 / 评论"),
        ("内容热度", "hot_score / view_count"),
        ("内容标签", "content_tag 多对多"),
    ]
    x = 0.75
    for title, sub in sources:
        add_card(slide, x, 1.25, 2.65, 0.95, title, sub, color=WHITE, accent=BLUE, icon=title[:2])
        connector(slide, x + 1.32, 2.2, 6.65, 3.0, color=GRAY_500, width=1)
        x += 3.05

    core = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.35), Inches(2.8), Inches(4.65), Inches(1.08))
    fill(core, BLUE, line=BLUE)
    add_text(slide, "综合得分 = 内容热度分 + 用户标签匹配分", 4.62, 3.12, 4.1, 0.28, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    formulas = [
        ("线路", "收藏×2 + 浏览 + 使用×3"),
        ("游记", "点赞×2 + 收藏×3 + 浏览 + 评论×2 + 偏好"),
        ("景点/美食", "hot_score×2 + 评分 + 标签偏好"),
    ]
    x = 1.05
    for title, formula in formulas:
        add_card(slide, x, 4.45, 3.45, 1.0, title, formula, color=WHITE, accent=GREEN)
        x += 4.02

    pill(slide, "Redis 热门结果缓存 30 分钟", 3.95, 5.95, 2.8, color=GREEN_LIGHT, line=GREEN, txt=GREEN)
    pill(slide, "Top N 降序输出", 7.0, 5.95, 2.0, color=BLUE_LIGHT, line=BLUE, txt=BLUE_DARK)
    add_footer_note(slide, "讲解重点：说明推荐不是空泛 AI，而是基于标签权重、行为数据和内容热度的可解释排序。")


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, "系统整体架构设计", 5)

    add_text(slide, "用户角色", 0.8, 1.0, 1.0, 0.25, size=13, color=GRAY_500, bold=True)
    add_card(slide, 0.65, 1.45, 1.5, 0.8, "游客/用户", "浏览、规划、互动", color=BLUE_LIGHT, accent=BLUE)
    add_card(slide, 0.65, 3.0, 1.5, 0.8, "管理员", "审核、配置、维护", color=GREEN_LIGHT, accent=GREEN)

    add_text(slide, "客户端", 2.65, 1.0, 1.0, 0.25, size=13, color=GRAY_500, bold=True)
    add_card(slide, 2.55, 1.45, 2.0, 0.8, "微信小程序", "UniApp + Vue3", color=WHITE, accent=BLUE)
    add_card(slide, 2.55, 3.0, 2.0, 0.8, "管理后台", "Vue3 + Element Plus", color=WHITE, accent=GREEN)

    add_text(slide, "服务层", 5.2, 1.0, 1.0, 0.25, size=13, color=GRAY_500, bold=True)
    add_card(slide, 5.05, 1.75, 2.0, 1.18, "REST API", "/api/v1/**\n统一 JSON 响应", color=BLUE_LIGHT, accent=BLUE)
    modules = ["user", "content", "route", "recommend", "travel", "common"]
    x = 7.35
    y = 1.28
    for i, m in enumerate(modules):
        pill(slide, m, x, y, 1.15, color=WHITE, line=GRAY_200, txt=GRAY_700)
        x += 1.25
        if i == 2:
            x = 7.35
            y += 0.62
    add_card(slide, 7.25, 2.65, 4.15, 1.0, "业务服务", "参数校验、事务处理、推荐计算、AI 文案编排、文件访问", color=WHITE, accent=BLUE_DARK)

    add_text(slide, "数据与外部能力", 5.2, 4.55, 1.8, 0.25, size=13, color=GRAY_500, bold=True)
    data = [("MySQL", "业务数据持久化"), ("Redis", "热门推荐缓存"), ("Uploads", "图片资源访问"), ("DeepSeek", "行程文案生成")]
    x = 2.1
    for title, sub in data:
        add_card(slide, x, 5.05, 2.0, 0.86, title, sub, color=WHITE, accent=ORANGE)
        x += 2.35

    for y in [1.85, 3.4]:
        connector(slide, 2.15, y, 2.5, y, color=GRAY_500, width=1.2)
        connector(slide, 4.55, y, 5.0, 2.34, color=GRAY_500, width=1.2)
    connector(slide, 7.05, 2.34, 7.25, 2.34, color=BLUE, width=1.5)
    connector(slide, 7.2, 3.65, 7.2, 5.0, color=GRAY_500, width=1.2)
    add_footer_note(slide, "讲解重点：从角色到数据流讲清楚，不需要展开每个接口。")


def slide_modules(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GRAY_100
    add_header(slide, "系统功能模块划分", 6)

    groups = [
        ("用户端功能", BLUE, [
            ("首页推荐", "展示热门路线、游记、景点、美食"),
            ("智能规划", "生成多日行程并查看详情"),
            ("旅游内容", "城市、景点、美食浏览与搜索"),
            ("社区互动", "游记发布、点赞、收藏、评论"),
            ("打卡足迹", "记录景点/美食打卡与轨迹"),
            ("个人中心", "资料、关注、消息与偏好标签"),
        ]),
        ("管理端功能", GREEN, [
            ("数据总览", "汇总用户、内容、访问等运营指标"),
            ("内容管理", "城市、景点、美食 CRUD"),
            ("游记审核", "处理用户发布内容的审核状态"),
            ("路线管理", "维护模板路线与行程明细"),
            ("运营配置", "活动专题、推荐策略、等级规则"),
            ("系统管理", "参数配置与操作日志"),
        ]),
        ("后端支撑", ORANGE, [
            ("接口服务", "RESTful API 统一对外"),
            ("推荐计算", "热度分与兴趣标签加权排序"),
            ("AI 能力", "DeepSeek 文案生成与兜底策略"),
            ("数据访问", "MyBatis Mapper 管理 SQL"),
            ("上传访问", "图片上传与静态资源映射"),
            ("异常处理", "统一异常与响应格式"),
        ]),
    ]

    x = 0.6
    for title, accent, items in groups:
        add_text(slide, title, x + 0.1, 1.12, 2.2, 0.28, size=15, color=accent, bold=True)
        y = 1.55
        for name, desc in items:
            add_card(slide, x, y, 3.85, 0.68, name, desc, color=WHITE, accent=accent)
            y += 0.78
        x += 4.25
    add_footer_note(slide, "讲解重点：按用户端、管理端、后端支撑三条线说明系统完整性。")


def slide_database(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, "数据库设计", 7)

    tables = {
        "user": (1.0, 1.3, BLUE),
        "user_tag": (2.65, 1.3, BLUE),
        "tag": (4.3, 1.3, BLUE),
        "city": (1.0, 3.1, GREEN),
        "scenic_spot": (2.65, 2.55, GREEN),
        "food": (2.65, 3.65, GREEN),
        "content_tag": (4.3, 3.1, GREEN),
        "travel_route": (6.1, 1.3, ORANGE),
        "route_day": (7.9, 1.3, ORANGE),
        "route_poi": (9.7, 1.3, ORANGE),
        "travel_note": (6.1, 3.4, CYAN),
        "comment": (7.9, 3.4, CYAN),
        "checkin_record": (9.7, 3.4, CYAN),
        "user_behavior": (6.1, 5.05, RED),
    }
    for name, (x, y, color) in tables.items():
        add_card(slide, x, y, 1.35, 0.52, name, None, color=WHITE, accent=color)
    relations = [
        ("user", "user_tag"), ("user_tag", "tag"), ("city", "scenic_spot"),
        ("city", "food"), ("scenic_spot", "content_tag"), ("food", "content_tag"),
        ("content_tag", "tag"), ("travel_route", "route_day"), ("route_day", "route_poi"),
        ("travel_note", "comment"), ("travel_note", "checkin_record"), ("user", "user_behavior"),
    ]
    for a, b in relations:
        ax, ay, _ = tables[a]
        bx, by, _ = tables[b]
        connector(slide, ax + 1.35, ay + 0.26, bx, by + 0.26, color=GRAY_500, width=1)

    add_card(
        slide, 0.9, 5.75, 4.25, 0.9, "关键字段",
        "user.openid、tag.weight、scenic_spot.hot_score、travel_route.ai_used、travel_note.status、user_behavior.score",
        color=BLUE_LIGHT, accent=BLUE,
    )
    add_card(
        slide, 5.55, 5.75, 5.85, 0.9, "设计说明",
        "行程三级结构适合多日展示；标签关联支持推荐与筛选；行为表沉淀个性化数据；del_flag 与时间字段便于审计。",
        color=GREEN_LIGHT, accent=GREEN,
    )
    add_footer_note(slide, "讲解重点：讲清楚表之间的业务关系，避免逐字段念表。")


def phone_mock(slide, x, y, title, lines, accent=BLUE):
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.1), Inches(4.0))
    fill(body, WHITE, line=GRAY_200)
    top = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.12), Inches(y + 0.14), Inches(1.86), Inches(0.54))
    fill(top, accent, line=accent)
    add_text(slide, title, x + 0.24, y + 0.32, 1.6, 0.16, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    yy = y + 0.92
    for line in lines:
        add_card(slide, x + 0.2, yy, 1.7, 0.52, line[0], line[1], color=GRAY_100, accent=accent)
        yy += 0.68


def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GRAY_100
    add_header(slide, "项目成果展示", 8)

    groups = [
        ("小程序首页", BLUE, [
            ("首页推荐", "路线、游记、景点与美食"),
            ("城市入口", "目的地切换与资源聚合"),
            ("热门内容", "个性化排序与快捷浏览"),
            ("底部导航", "首页、行程、社区、我的"),
        ]),
        ("智能规划", GREEN, [
            ("选择城市", "天数、预算、同行人设置"),
            ("兴趣标签", "历史、自然、美食等偏好"),
            ("必选 POI", "景点与餐饮按天编排"),
            ("生成行程", "早中晚安排与详情查看"),
        ]),
        ("管理后台", ORANGE, [
            ("数据总览", "用户、内容、访问等指标"),
            ("内容管理", "城市、景点、美食 CRUD"),
            ("游记审核", "用户内容通过或驳回"),
            ("策略配置", "推荐策略与等级规则"),
        ]),
    ]

    x = 0.62
    for title, accent, items in groups:
        add_text(slide, title, x + 0.18, 1.1, 2.1, 0.28, size=15, color=accent, bold=True)
        y = 1.55
        for name, desc in items:
            add_card(slide, x, y, 3.82, 0.78, name, desc, color=WHITE, accent=accent)
            y += 0.9
        x += 4.25

    add_text(slide, "项目亮点", 0.75, 5.45, 1.2, 0.25, size=14, color=BLUE_DARK, bold=True)
    highlights = [
        "三端联调完整",
        "AI 文案与规则兜底",
        "标签权重推荐",
        "UGC 互动闭环",
        "后台运营配置",
    ]
    x = 2.0
    for item in highlights:
        pill(slide, item, x, 5.38, 1.55, color=BLUE_LIGHT, line=BLUE, txt=BLUE_DARK)
        x += 1.75

    add_text(slide, "演示流程：登录/选择兴趣 → 生成智能行程 → 查看详情并打卡 → 发布游记 → 后台审核", 0.75, 6.22, 9.5, 0.26, size=11, color=GRAY_700)
    add_text(slide, "感谢各位老师聆听，敬请批评指正！", 8.7, 6.65, 3.4, 0.35, size=16, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "08/08", 12.1, 6.95, 0.7, 0.25, size=9, color=GRAY_500, align=PP_ALIGN.RIGHT)


def write_outline():
    lines = [
        "# 智能旅游系统毕业答辩 PPT 内容与讲稿",
        "",
        "> 总页数：8 页；自述时长建议：8 分钟，其中系统演示 3 到 4 分钟。",
        "",
    ]
    for idx, info in enumerate(SLIDES, start=1):
        lines.extend([
            f"## 第 {idx} 页：{info['title']}",
            "",
            f"- 页面布局建议：{info['layout']}",
            f"- 配图/图示建议：{info['visual']}",
            f"- 答辩讲稿：{info['speech']}",
            "",
        ])
    lines.extend([
        "## 完整 PPT 页结构",
        "",
        "1. 封面：毕业设计题目、学生姓名、班级、学校名称",
        "2. 项目开发环境与技术栈",
        "3. 核心算法（一）：智能行程规划",
        "4. 核心算法（二）：个性化推荐",
        "5. 系统整体架构设计",
        "6. 系统功能模块划分",
        "7. 数据库设计",
        "8. 项目成果展示与致谢",
        "",
    ])
    OUTLINE_PATH.write_text("\n".join(lines), encoding="utf-8")


def build():
    prs = Presentation()
    prs.slide_width = WIDE
    prs.slide_height = HIGH
    slide_cover(prs)
    slide_tech(prs)
    slide_algorithm_route(prs)
    slide_algorithm_recommend(prs)
    slide_architecture(prs)
    slide_modules(prs)
    slide_database(prs)
    slide_results(prs)
    prs.save(OUTPUT_PATH)
    write_outline()
    print(f"Generated: {OUTPUT_PATH}")
    print(f"Generated: {OUTLINE_PATH}")


if __name__ == "__main__":
    build()
